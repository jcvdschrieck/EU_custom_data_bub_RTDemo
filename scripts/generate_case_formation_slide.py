#!/usr/bin/env python3
"""
Generate a single-slide PowerPoint explaining how transactions are
grouped into a case: the Jaccard-similarity linking rule, the
append-while-open policy, and how case-level fields (name, risk score,
engine scores, AI rationale) are derived from the linked orders.

Output: docs/case_formation.pptx

Uses the same palette and helpers as scripts/generate_risk_engines_deck.py
so the slide drops cleanly into the existing deck.
"""
from pathlib import Path

from pptx             import Presentation
from pptx.dml.color   import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text   import PP_ALIGN, MSO_ANCHOR
from pptx.util        import Inches, Pt

# ── Palette (matches generate_risk_engines_deck.py) ────────────────────────
EU_BLUE       = RGBColor(0x00, 0x33, 0x99)
EU_GOLD       = RGBColor(0xFF, 0xCC, 0x00)
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)
LIGHT_TEXT    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GREEN  = RGBColor(0x1F, 0x7A, 0x3C)
ACCENT_AMBER  = RGBColor(0xE6, 0x82, 0x0A)
ACCENT_RED    = RGBColor(0xC0, 0x39, 0x2B)
SUBTLE_GREY   = RGBColor(0x66, 0x66, 0x66)
BG_LIGHT      = RGBColor(0xF5, 0xF5, 0xF5)
BOX_BG        = RGBColor(0xFF, 0xFF, 0xFF)
HILITE_BG     = RGBColor(0xFF, 0xF3, 0xE0)
HILITE_EDGE   = RGBColor(0xE6, 0x82, 0x0A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────

def _add_rect(slide, left, top, width, height, *, fill, edge=None,
              text=None, font_size=12, bold=False, color=DARK_TEXT,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1.25)
    if text is not None:
        tf = shape.text_frame
        tf.vertical_anchor = anchor
        tf.margin_left  = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return shape


def _add_text_box(slide, left, top, width, height, text, *,
                  font_size=12, bold=False, color=DARK_TEXT,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  fill=None, edge=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left  = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top   = Inches(0.05)
    tf.margin_bottom= Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return shape


def _add_bullet_box(slide, left, top, width, height, title, bullets, *,
                    title_size=14, body_size=11, fill=BOX_BG, edge=EU_BLUE,
                    title_color=EU_BLUE, body_color=DARK_TEXT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = edge
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left   = Inches(0.12)
    tf.margin_right  = Inches(0.12)
    tf.margin_top    = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    run = p0.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(title_size)
    run.font.color.rgb = title_color
    run.font.name = "Calibri"

    for b in bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(body_size)
        run.font.color.rgb = body_color
        run.font.name = "Calibri"


def _add_arrow(slide, x1, y1, x2, y2, color=SUBTLE_GREY):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


# ── Slide builder ──────────────────────────────────────────────────────────

def slide_case_formation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    _add_text_box(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55),
                  "How orders are grouped into a case",
                  font_size=28, bold=True, color=EU_BLUE)
    _add_text_box(s, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.4),
                  "Every investigate-routed transaction either joins an "
                  "existing open case or starts a new one. The decision "
                  "is deterministic and reversible.",
                  font_size=13, color=SUBTLE_GREY)

    # ── Linking rule (left) ──────────────────────────────────────────────
    _add_bullet_box(s,
        Inches(0.5), Inches(1.4), Inches(6.3), Inches(2.6),
        "Linking rule — find_similar_open_case",
        [
            "Applied when a transaction hits the Custom & Tax Risk "
            "Management factory and has been routed to investigate.",
            "Among all cases whose Status is not Closed, look for one "
            "that shares, with the incoming order, the same Seller, "
            "the same declared Product Category, and the same Country "
            "of Destination.",
            "For the candidates that survive that filter, compare their "
            "primary Product Description to the new order's description "
            "using Jaccard similarity — keep the best match above the "
            "0.40 threshold.",
            "Match → append the order to that existing case. No match "
            "→ create a new case with this order as the primary.",
        ],
    )

    # ── Jaccard similarity box (right) ──────────────────────────────────
    _add_bullet_box(s,
        Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.6),
        "Jaccard similarity in plain terms",
        [
            "Take the two product descriptions, lower-case them, split "
            "each into a set of unique words.",
            "similarity = |A ∩ B| / |A ∪ B|  — the fraction of unique "
            "words the two descriptions share.",
            "Threshold 0.40: two descriptions must share at least 40 % "
            "of their combined vocabulary to be treated as the same "
            "product cluster.",
            "Example match: \"Bluetooth wireless earbuds\" vs "
            "\"Bluetooth wireless earphones\" → 0.50 → same case.",
            "Example mismatch: \"Bluetooth wireless earbuds\" vs "
            "\"4K webcam HDR\" → 0.00 → separate cases.",
        ],
    )

    # ── Append-while-open policy + case derivation (bottom row) ─────────

    # Append-while-open policy
    _add_bullet_box(s,
        Inches(0.5), Inches(4.15), Inches(6.3), Inches(2.95),
        "Append-while-open policy",
        [
            "A case acts as a live container: as long as it is not "
            "Closed, every new order that passes the linking rule is "
            "appended to it.",
            "When the case closes (retainment, release, refused, or the "
            "agent wraps up), the container is sealed. A new matching "
            "order arriving afterwards starts a fresh case — the sealed "
            "one is not reopened.",
            "On each append the per-engine case scores are re-averaged "
            "across the linked orders and the Overall_Case_Risk_Score "
            "is recomputed. The case's AI recommendation is re-derived "
            "on every read, so it stays consistent with the new order "
            "mix without requiring any manual refresh.",
            "Keeps cases coherent (same seller / category / destination "
            "/ description cluster) without unbounded growth — an "
            "unrelated new pattern cannot contaminate an existing case.",
        ],
    )

    # Case derivation
    _add_bullet_box(s,
        Inches(7.0), Inches(4.15), Inches(5.8), Inches(2.95),
        "Case-level fields — where they come from",
        [
            "Case name: built from the unique product descriptions of "
            "the linked orders (top two, plus \"+N more\" if longer), "
            "followed by the seller name. Composed in the frontend by "
            "backendCaseToCase.",
            "Seller, origin, destination, declared category: taken "
            "from the case's primary order (first linked Sales_Order).",
            "Engine scores on the case (VAT ratio, ML, IE watchlist, "
            "vagueness): averaged across the linked orders at append "
            "time. Overall_Case_Risk_Score is their weighted sum, "
            "capped at 1.0.",
            "AI-suggested action + rationale: computed at hydration "
            "time by _compute_customs_recommendation — historical "
            "retain-ratio for this (seller, category, destination) "
            "plus any confirming engine signals on the current case.",
            "VAT_Gap_Fee: populated only once the VAT Fraud Detection "
            "agent returns a definitive verdict (sum of per-order "
            "(expected − declared) VAT across linked orders).",
        ],
    )


# ── Main ───────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_case_formation(prs)

    out = Path(__file__).resolve().parent.parent / "docs" / "case_formation.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(out)
    print(f"Wrote {out} ({out.stat().st_size:,} bytes, {len(prs.slides)} slide)")


if __name__ == "__main__":
    main()
