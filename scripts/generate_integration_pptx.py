"""
Generate a 4-slide architecture deck for the EU Customs Data Hub +
Revenue Guardian integration.

Slides:
  1. Component architecture (frontend + backend, info gateways)
  2. Data architecture
  3. Functional architecture + HL integration changes
  4. Backend/Frontend vs EU_custom_data_hub/revenue-guardian split
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme ────────────────────────────────────────────────────────────────────
EU_BLUE      = RGBColor(0x00, 0x3D, 0xA5)
EU_BLUE_LT   = RGBColor(0xE3, 0xEC, 0xF8)
PURPLE       = RGBColor(0x6F, 0x42, 0xC1)
PURPLE_LT    = RGBColor(0xEF, 0xE7, 0xFB)
GREEN        = RGBColor(0x19, 0x87, 0x54)
GREEN_LT     = RGBColor(0xE6, 0xF4, 0xEA)
RED          = RGBColor(0xDC, 0x35, 0x45)
RED_LT       = RGBColor(0xFB, 0xEA, 0xEC)
AMBER        = RGBColor(0xFD, 0x7E, 0x14)
AMBER_LT     = RGBColor(0xFF, 0xF4, 0xE6)
GREY         = RGBColor(0x49, 0x50, 0x57)
GREY_LT      = RGBColor(0xF1, 0xF3, 0xF5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEAL         = RGBColor(0x0B, 0x7A, 0x75)
TEAL_LT      = RGBColor(0xDC, 0xF1, 0xEF)


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = EU_BLUE


def add_subtitle(slide, text, top=0.78):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(12.5), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = GREY


def add_box(slide, x, y, w, h, label, fill, border, font_size=10,
            bold=True, font_color=None, sub=None, sub_size=8):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color or border
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = sub
        p2.font.size = Pt(sub_size)
        p2.font.bold = False
        p2.font.color.rgb = GREY
    return shape


def add_zone(slide, x, y, w, h, label, fill, border, label_color=None):
    """Faint background zone with a small label badge top-left."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.line.dash_style = 7  # dash
    # label badge
    badge = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.05),
                                     Inches(w - 0.1), Inches(0.25))
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = label_color or border
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=GREY, width=1.5, label=None,
              dashed=False, label_offset=(0, -0.15)):
    """Straight connector line with optional label and dash style."""
    line = slide.shapes.add_connector(2, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = 7
    # arrow head
    line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = line.line._get_or_add_ln()
    tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("h", "med")
    if label:
        midx = (x1 + x2) / 2 + label_offset[0]
        midy = (y1 + y2) / 2 + label_offset[1]
        box = slide.shapes.add_textbox(Inches(midx - 0.9), Inches(midy),
                                       Inches(1.8), Inches(0.22))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = color


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5),
                                   Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = GREY


def add_bullets(slide, x, y, w, h, bullets, title=None,
                fill=WHITE, border=EU_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = border
        first = False
    else:
        first = True
    for i, b in enumerate(bullets):
        if first and i == 0:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(9)
        p.font.color.rgb = GREY


# ── Build presentation ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ============================================================================
# Slide 1 — Component architecture + information gateways
# ============================================================================
s = prs.slides.add_slide(blank)
add_title(s, "1. Component architecture — Frontend + Backend & information gateways")
add_subtitle(s, "Two independent processes (port 8000 & 8080). All cross-process traffic goes through one HTTP API surface.")

# Backend zone (left)
add_zone(s, 0.4, 1.2, 6.8, 5.6, "BACKEND  •  EU_custom_data_hub_RTDemo  •  FastAPI  •  port 8000", EU_BLUE_LT, EU_BLUE)
# Frontend zone (right)
add_zone(s, 7.45, 1.2, 5.5, 5.6, "FRONTEND  •  revenue-guardian  •  Vite + React  •  port 8080", PURPLE_LT, PURPLE)

# Backend blocks
add_box(s, 0.7, 1.7, 1.9, 0.55, "Simulation engine", WHITE, EU_BLUE, 9, sub="replay of March-2026 txs", sub_size=7)
add_box(s, 0.7, 2.45, 1.9, 0.55, "Risk engines", WHITE, EU_BLUE, 9, sub="vat_ratio + watchlist", sub_size=7)
add_box(s, 2.85, 2.45, 2.0, 0.55, "Release Factory", WHITE, EU_BLUE, 9, sub="route: release/retain/investigate", sub_size=7)
add_box(s, 5.1, 2.45, 1.95, 0.55, "C&T Risk Mgmt Factory", WHITE, PURPLE, 9, sub="rewritten — opens cases", sub_size=7)
add_box(s, 5.1, 3.25, 1.95, 0.55, "DB Store Factory", WHITE, EU_BLUE, 9, sub="emits Custom Outcome", sub_size=7)
add_box(s, 0.7, 3.25, 1.9, 0.55, "Pub/sub broker", WHITE, EU_BLUE, 9, sub="in-memory (lib/broker)", sub_size=7)
add_box(s, 2.85, 3.25, 2.0, 0.55, "Event store", WHITE, EU_BLUE, 9, sub="data/events/*.json", sub_size=7)

add_box(s, 0.7, 4.05, 2.0, 0.6, "investigation.db", GREEN_LT, GREEN, 10, sub="3 tables (case + order + risk)", sub_size=7)
add_box(s, 2.85, 4.05, 2.0, 0.6, "european_custom.db", GREY_LT, GREY, 10, sub="legacy — sink deactivated", sub_size=7)
add_box(s, 5.1, 4.05, 1.95, 0.6, "VAT Fraud Agent", GREY_LT, GREY, 10, sub="LM Studio bridge — not yet wired", sub_size=7)

# REST/SSE surface
add_box(s, 0.7, 5.0, 6.35, 0.65,
        "REST + SSE API surface  /api/rg/cases/*",
        TEAL_LT, TEAL, 11,
        sub="GET cases · GET cases/stream · POST customs-action · tax-action · final-decision · communication",
        sub_size=8)

# Built-in simulation UI (still on the backend!)
add_box(s, 0.7, 5.85, 6.35, 0.7,
        "Simulation / pipeline diagram UI",
        EU_BLUE_LT, EU_BLUE, 11,
        sub="React SPA built into the backend repo, served by FastAPI at :8000",
        sub_size=8)

# Frontend blocks
add_box(s, 7.7, 1.7, 4.95, 0.55, "App.tsx  •  router  •  bootstraps caseStore", WHITE, PURPLE, 9)
add_box(s, 7.7, 2.4, 2.4, 0.6, "apiClient.ts", PURPLE_LT, PURPLE, 10,
        sub="REST + SSE + action POSTs", sub_size=7)
add_box(s, 10.25, 2.4, 2.4, 0.6, "caseStore.ts", PURPLE_LT, PURPLE, 10,
        sub="in-memory cases cache", sub_size=7)
add_box(s, 7.7, 3.2, 4.95, 0.55, "taxReviewStore.ts (modified)", WHITE, PURPLE, 9,
        sub="reads from caseStore · setters POST through", sub_size=7)

add_box(s, 7.7, 3.95, 1.55, 0.7, "Customs Authority", WHITE, PURPLE, 9, sub="page (modified)", sub_size=7)
add_box(s, 9.4, 3.95, 1.55, 0.7, "Tax Authority", WHITE, PURPLE, 9, sub="page", sub_size=7)
add_box(s, 11.1, 3.95, 1.55, 0.7, "Investigations", WHITE, PURPLE, 9, sub="page", sub_size=7)

add_box(s, 7.7, 4.85, 4.95, 0.55, "TaxReviewDialog · EucdhTransactionDialog · LaunchInvestigationDialog",
        WHITE, PURPLE, 9, sub="native shadcn/ui dialogs, unchanged", sub_size=7)

add_box(s, 7.7, 5.55, 4.95, 0.55, "localStorage  (UI cache only)", GREY_LT, GREY, 9,
        sub="recommendation pick · per-officer flags · risk breakdown overrides",
        sub_size=7)

# Gateway arrows (the two info gateways!)
# REST gateway
add_arrow(s, 7.05, 5.3, 7.7, 2.7, color=TEAL, width=2.5,
          label="REST  /api/rg/*", label_offset=(0.15, -0.25))
# SSE gateway
add_arrow(s, 7.05, 5.5, 7.7, 2.95, color=TEAL, width=2.5, dashed=True,
          label="SSE  cases/stream", label_offset=(0.25, 0.05))
# Action POST back
add_arrow(s, 7.7, 3.5, 7.05, 5.4, color=AMBER, width=1.75,
          label="POST  actions", label_offset=(-0.6, -0.05))

add_footer(s, "Information gateways: 1) REST GET /api/rg/cases  2) SSE /api/rg/cases/stream (push)  3) REST POST officer actions back. CORS open. No shared state.")

# ============================================================================
# Slide 2 — Data architecture
# ============================================================================
s = prs.slides.add_slide(blank)
add_title(s, "2. Data architecture")
add_subtitle(s, "Three SQLite stores, two of which are now decoupled. Snapshot semantics on case-open.")

# investigation.db (active, 3 tables)
add_zone(s, 0.4, 1.2, 6.8, 4.0, "investigation.db  •  authoritative for Revenue Guardian", PURPLE_LT, PURPLE)

add_box(s, 0.7, 1.7, 2.0, 1.55,
        "Sales_Order\n(snapshot)",
        WHITE, PURPLE, 10,
        sub="ID · BusinessKey · HS_Category · Description · Value · VAT_Rate · VAT_Fee · Seller · Origin · Destination · Status · Update_time",
        sub_size=7)
add_box(s, 2.85, 1.7, 2.0, 1.55,
        "Sales_Order_Risk\n(snapshot)",
        WHITE, PURPLE, 10,
        sub="Risk_Type · Overall_Score · Overall_Level · Confidence · Proposed_Action · 4 dimensional sub-scores (null for now)",
        sub_size=7)
add_box(s, 5.0, 1.7, 2.0, 1.55,
        "Sales_Order_Case\n(workspace)",
        PURPLE_LT, PURPLE, 10,
        sub="Status · VAT_Problem_Type · Recommended_* · AI_Analysis · VAT_Gap_Fee · Proposed_Action_Tax · Proposed_Action_Customs · Communication[]",
        sub_size=7)
# join arrow
add_arrow(s, 1.7, 3.35, 1.7, 3.7, color=PURPLE)
add_arrow(s, 3.85, 3.35, 3.85, 3.7, color=PURPLE)
add_arrow(s, 6.0, 3.35, 6.0, 3.7, color=PURPLE)
add_box(s, 0.7, 3.7, 6.3, 0.5, "joined on  Sales_Order_Business_Key  →  hydrated DTO returned by /api/rg/cases",
        WHITE, PURPLE, 10)

# legacy
add_zone(s, 0.4, 5.35, 6.8, 0.95, "european_custom.db  •  legacy (DB Store sink deactivated)",
         GREY_LT, GREY)
add_box(s, 0.7, 5.65, 2.0, 0.55, "transactions", WHITE, GREY, 9,
        sub="historical seed (risk-engine baselines)", sub_size=7)
add_box(s, 2.85, 5.65, 2.0, 0.55, "Sales_Order", WHITE, GREY, 9,
        sub="not written anymore", sub_size=7)
add_box(s, 5.0, 5.65, 2.0, 0.55, "Sales_Order_Risk", WHITE, GREY, 9,
        sub="not written anymore", sub_size=7)

# Reference / lookup tables — now hosted in the backend (european_custom.db)
add_zone(s, 0.4, 6.4, 6.8, 0.95,
         "Reference tables  •  european_custom.db  →  GET /api/reference  (single fetch on app startup)",
         AMBER_LT, AMBER)
add_box(s, 0.55, 6.65, 1.55, 0.6, "vat_categories",
        WHITE, AMBER, 9,
        sub="label · rate · description", sub_size=7)
add_box(s, 2.2, 6.65, 1.55, 0.6, "risk_levels",
        WHITE, AMBER, 9,
        sub="name · display_color", sub_size=7)
add_box(s, 3.85, 6.65, 1.55, 0.6, "eu_regions",
        WHITE, AMBER, 9,
        sub="country_code → region", sub_size=7)
add_box(s, 5.5, 6.65, 1.6, 0.6, "suspicion_types",
        WHITE, AMBER, 9,
        sub="name · description · icon · color", sub_size=7)

# event store
add_zone(s, 7.45, 1.2, 5.5, 4.0, "Event store  •  data/events/<topic>/", TEAL_LT, TEAL)

add_box(s, 7.75, 1.7, 5.0, 0.5, "sales_order_event", WHITE, TEAL, 9)
add_box(s, 7.75, 2.25, 5.0, 0.5, "rt_risk_outcome  ·  order_validation", WHITE, TEAL, 9)
add_box(s, 7.75, 2.8, 5.0, 0.5, "assessment_outcome  (route: release/retain/investigate)", WHITE, TEAL, 9)
add_box(s, 7.75, 3.35, 5.0, 0.5, "investigation_outcome  (case closure exit event)", WHITE, TEAL, 9)
add_box(s, 7.75, 3.9, 5.0, 0.7, "custom_outcome  (NEW terminal broker)",
        TEAL_LT, TEAL, 10,
        sub="status: automated_release · custom_release · custom_retain", sub_size=7)
add_box(s, 7.75, 4.7, 5.0, 0.4, "(events persisted as JSON for replay & counters)",
        WHITE, TEAL, 8, bold=False)

# Frontend cache
add_zone(s, 7.45, 5.35, 5.5, 2.0, "revenue-guardian frontend memory", PURPLE_LT, PURPLE)
add_box(s, 7.75, 5.85, 5.0, 0.45, "caseStore: Map<Case_ID, BackendCase>",
        WHITE, PURPLE, 10, sub="hydrated from REST + SSE", sub_size=7)
add_box(s, 7.75, 6.35, 5.0, 0.45, "referenceStore: lookup bundle",
        WHITE, PURPLE, 10, sub="hydrated once from /api/reference", sub_size=7)
add_box(s, 7.75, 6.85, 5.0, 0.5, "localStorage (UI cache, non-authoritative)",
        GREY_LT, GREY, 9, sub="recommendation pick · per-officer flags · risk overrides", sub_size=7)

add_footer(s, "Snapshot semantics: case-open freezes order + risk into investigation.db. Lookups now live in european_custom.db, fetched once via GET /api/reference at app startup.")

# ============================================================================
# Slide 3 — Functional architecture + integration changes (HL)
# ============================================================================
s = prs.slides.add_slide(blank)
add_title(s, "3. Functional architecture & integration changes (high level)")
add_subtitle(s, "Same functional pipeline as before; the C&T block was rewired and the data hub sink was replaced by a terminal event broker.")

# Pipeline strip
y = 1.25
strip_h = 3.4
add_zone(s, 0.4, y, 12.55, strip_h, "Functional pipeline (left → right)", EU_BLUE_LT, EU_BLUE)

# Boxes along the pipeline
bx_w, bx_h = 1.5, 0.85
top = y + 1.2
add_box(s, 0.55, top, bx_w, bx_h, "Entry", WHITE, GREY, 10, sub="seed orders", sub_size=7)
add_box(s, 2.2, top, bx_w, bx_h, "Risk Assessment", WHITE, EU_BLUE, 10, sub="2 engines + validation", sub_size=7)
add_box(s, 3.85, top, bx_w, bx_h, "Release Factory", WHITE, EU_BLUE, 10, sub="route by colour", sub_size=7)
add_box(s, 5.5, top, bx_w, bx_h, "C&T Risk Mgmt", PURPLE_LT, PURPLE, 10, sub="opens cases", sub_size=7)
add_box(s, 7.15, top, bx_w, bx_h, "Officer review", PURPLE_LT, PURPLE, 10, sub="Customs / Tax UI", sub_size=7)
add_box(s, 8.8, top, bx_w, bx_h, "DB Store", WHITE, EU_BLUE, 10, sub="emits terminal", sub_size=7)
add_box(s, 10.45, top, bx_w, bx_h, "Custom Outcome", TEAL_LT, TEAL, 10, sub="3 statuses", sub_size=7)

# arrows
ay = top + bx_h / 2
for x in (2.05, 3.7, 5.35, 7.0, 8.65, 10.3):
    add_arrow(s, x, ay, x + 0.15, ay, color=GREY)

# Convergence note (ties to Slide 1)
note = slide_box = s.shapes.add_textbox(Inches(0.55), Inches(top - 0.5),
                                         Inches(12.4), Inches(0.4))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "Convergence with Slide 1:  Risk Assessment → Release Factory → DB Store all live in the backend; C&T + Officer review live across the gateway (backend stores cases, frontend renders/acts)."
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = GREY

# Changes panel
ch_top = 4.85
add_bullets(s, 0.4, ch_top, 4.15, 2.45,
            title="Native — unchanged",
            fill=EU_BLUE_LT, border=EU_BLUE,
            bullets=[
                "Simulation engine + replay clock",
                "Risk engines (vat_ratio, watchlist) + validation",
                "Release Factory routing logic",
                "Pub/sub broker + event store mechanism",
                "Simulation / pipeline UI (built-in)",
                "VAT Fraud Detection agent bridge",
            ])
add_bullets(s, 4.65, ch_top, 4.15, 2.45,
            title="Modified — C&T integration",
            fill=PURPLE_LT, border=PURPLE,
            bullets=[
                "C&T Factory now opens cases (3-row write)",
                "INVESTIGATION_OUTCOME = case-closure event",
                "/api/rg/cases REST + SSE surface added",
                "Officer-action endpoints (close, comment, …)",
                "investigation.db extended to 3 tables",
                "Officer actions flow back over HTTP",
            ])
add_bullets(s, 8.9, ch_top, 4.05, 2.45,
            title="New — terminal broker",
            fill=TEAL_LT, border=TEAL,
            bullets=[
                "DB Store sink deactivated (no hub writes)",
                "DB Store now publishes Custom Outcome",
                "Per-status counters (auto / custom rls / retain)",
                "Pipeline diagram swaps the cylinder for a broker",
                "Counters surface on /api/simulation/pipeline",
                "Backwards-compat: european_custom.db still seeded",
            ])

add_footer(s, "All three columns map back to the components on Slide 1; Custom Outcome is the new exit point of the pipeline.")

# ============================================================================
# Slide 4 — Two viewpoints: backend/frontend  vs  EU_custom_data_hub/revenue-guardian
# ============================================================================
s = prs.slides.add_slide(blank)
add_title(s, "4. Two viewpoints — process vs. project")
add_subtitle(s, "These are not the same split. The backend repo also owns a built-in frontend (the simulation page).")

# Headers
def header(slide, x, y, w, text, color):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color

# Layout: 2 columns × 2 rows
# Top row = process role (Backend / Frontend)
# Left col = EU_custom_data_hub repo, Right col = revenue-guardian repo
header(s, 4.4, 0.95, 4.0, "Backend (server, port 8000)", EU_BLUE)
header(s, 8.55, 0.95, 4.0, "Frontend (UI, port 8080)", PURPLE)

# Vertical labels for the row dimension
def vlabel(slide, x, y, h, text, color):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.7), Inches(h))
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color

# Two rows, two columns
COL1_X = 4.4   # Backend column
COL2_X = 8.55  # Frontend column
ROW1_Y = 1.4
ROW2_Y = 4.45
ROW_H  = 2.85

# Row 1 — EU_custom_data_hub repo
add_zone(s, 0.4, ROW1_Y, 12.55, ROW_H,
         "EU_custom_data_hub_RTDemo  (repo)", EU_BLUE_LT, EU_BLUE)
vlabel(s, 0.55, ROW1_Y + 0.4, ROW_H - 0.5, "EU_custom_data_hub  REPO", EU_BLUE)

add_bullets(s, COL1_X, ROW1_Y + 0.4, 4.0, ROW_H - 0.55,
            title="Backend role (FastAPI)",
            fill=WHITE, border=EU_BLUE,
            bullets=[
                "Simulation engine + clock",
                "Pub/sub broker + event store",
                "Risk engines, Release Factory, C&T Factory",
                "DB Store Factory → Custom Outcome",
                "investigation.db (3 tables) — case data",
                "REST + SSE: /api/rg/cases/*",
                "Officer-action endpoints",
                "VAT Fraud Detection agent bridge",
            ])

add_bullets(s, COL2_X, ROW1_Y + 0.4, 4.0, ROW_H - 0.55,
            title="Frontend role  (this repo also!)",
            fill=PURPLE_LT, border=PURPLE,
            bullets=[
                "Simulation / pipeline UI — React SPA",
                "Built into frontend/dist, served by FastAPI",
                "Operator audience: demo controller, not officer",
                "URL: http://localhost:8000",
                "Pipeline diagram with broker + factory tiles",
                "Custom Outcome counters live here",
                "(Independent of revenue-guardian)",
            ])

# Row 2 — revenue-guardian repo
add_zone(s, 0.4, ROW2_Y, 12.55, ROW_H,
         "revenue-guardian  (repo)", PURPLE_LT, PURPLE)
vlabel(s, 0.55, ROW2_Y + 0.4, ROW_H - 0.5, "revenue-guardian  REPO", PURPLE)

add_bullets(s, COL1_X, ROW2_Y + 0.4, 4.0, ROW_H - 0.55,
            title="Backend role  —  none",
            fill=GREY_LT, border=GREY,
            bullets=[
                "No server in this repo",
                "All persistence delegated to backend",
                "All authoritative state lives at :8000",
                "(The .env points to VITE_API_BASE_URL)",
            ])

add_bullets(s, COL2_X, ROW2_Y + 0.4, 4.0, ROW_H - 0.55,
            title="Frontend role (Vite + React, port 8080)",
            fill=WHITE, border=PURPLE,
            bullets=[
                "Customs Authority dashboard (officer)",
                "Tax Authority dashboard (officer)",
                "Investigations review page",
                "apiClient + caseStore (REST + SSE)",
                "taxReviewStore (officer-action POSTs)",
                "URL: http://localhost:8080",
                "Audience: customs / tax officers",
            ])

add_footer(s,
    "Key insight: the backend repo carries its OWN frontend (simulation page). The 'frontend repo' is only one of two UIs. "
    "Process split (port 8000 vs 8080) is orthogonal to repo split (data-hub repo vs guardian repo).")


# ── Save ────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent.parent / "docs" / "EU_Custom_DataHub_Integration.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"Wrote {out}")
