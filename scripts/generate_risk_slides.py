#!/usr/bin/env python3
"""
Generate a PowerPoint deck documenting the two risk monitoring rules
and the consolidation/scoring mechanism.

Output: docs/risk_monitoring_rules.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Palette ────────────────────────────────────────────────────────────────

EU_BLUE      = RGBColor(0x00, 0x33, 0x99)
EU_GOLD      = RGBColor(0xFF, 0xCC, 0x00)
DARK_TEXT     = RGBColor(0x22, 0x22, 0x22)
LIGHT_TEXT    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GREEN  = RGBColor(0x1F, 0x7A, 0x3C)
ACCENT_AMBER  = RGBColor(0xE6, 0x82, 0x0A)
ACCENT_RED    = RGBColor(0xC0, 0x39, 0x2B)
SUBTLE_GREY   = RGBColor(0x66, 0x66, 0x66)
BG_LIGHT      = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER  = RGBColor(0x00, 0x50, 0xA0)
TABLE_ROW_ALT = RGBColor(0xE8, 0xF0, 0xFA)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)


def _set_cell(cell, text, font_size=11, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Blue banner
    from pptx.util import Emu
    shp = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(2.8))
    shp.fill.solid()
    shp.fill.fore_color.rgb = EU_BLUE
    shp.line.fill.background()
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = LIGHT_TEXT
    # Subtitle
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = EU_GOLD
    return slide


def _add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Thin blue bar at top
    shp = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    shp.fill.solid()
    shp.fill.fore_color.rgb = EU_BLUE
    shp.line.fill.background()
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = LIGHT_TEXT
    return slide


def _add_bullet_box(slide, left, top, width, height, bullets, font_size=12):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)
    return txBox


def _add_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a styled table. rows_data = list of lists, first row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top),
        Inches(width), Inches(0.35 * n_rows))
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            is_header = (r == 0)
            _set_cell(cell, str(val),
                       font_size=10 if not is_header else 11,
                       bold=is_header,
                       color=LIGHT_TEXT if is_header else DARK_TEXT,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
    return tbl_shape


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ── Slide 1: Title ──────────────────────────────────────────────────
    _add_title_slide(prs,
        "Risk Monitoring Rules",
        "EU Custom Data Hub — Real-Time Risk Assessment")

    # ── Slide 2: RT Risk Monitoring 1 — VAT Ratio Deviation ────────────
    s2 = _add_content_slide(prs, "RT Risk Monitoring 1 — VAT Ratio Deviation")

    _add_bullet_box(s2, 0.5, 1.1, 5.5, 2.5, [
        ("Engine ID: vat_ratio", 0),
        ("Detects sudden shifts in a supplier's VAT-to-value ratio for a specific destination country", 0),
        ("", 0),
        ("Algorithm:", 0),
        ("Compute VAT/value ratio for (seller, buyer_country) over the last 7 days", 1),
        ("Compare to the same ratio over the preceding 8 weeks (days −63 to −7)", 1),
        ("If deviation > 25% of the historical ratio → raise alarm (7-day expiry)", 1),
        ("While alarm is active, tag all new transactions from the same pair as suspicious", 1),
    ], font_size=11)

    _add_table(s2, 0.5, 3.8, 5.5, [
        ["Parameter", "Value", "Description"],
        ["MIN_CURRENT_TX", "3", "Min transactions in 7-day window"],
        ["MIN_HISTORICAL_TX", "5", "Min transactions in 8-week baseline"],
        ["DEVIATION_THRESHOLD", "25%", "Trigger threshold"],
        ["SUSPICIOUS_COUNTRIES", "{IE}", "Countries entering suspicious queue"],
    ], col_widths=[2.2, 1.0, 2.3])

    # Scenario box on the right
    txBox = s2.shapes.add_textbox(Inches(6.3), Inches(1.1), Inches(3.3), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Seeded Scenario"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = EU_BLUE
    for text in [
        "",
        "TechZone GmbH (SUP001, DE) → Ireland",
        "",
        "Week 2 of March 2026: electronics billed at 0% instead of the correct 23%.",
        "",
        "7-day VAT ratio drops from ~19% to ~0%, far exceeding the 25% threshold.",
        "",
        "Alarm fires on the first affected transaction.",
    ]:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = text
        run2.font.size = Pt(10)
        run2.font.color.rgb = DARK_TEXT if text else SUBTLE_GREY
        p2.space_after = Pt(2)

    # ── Slide 3: RT Risk Monitoring 2 — Watchlist ──────────────────────
    s3 = _add_content_slide(prs, "RT Risk Monitoring 2 — Supplier/Origin Watchlist")

    _add_bullet_box(s3, 0.5, 1.1, 5.5, 2.0, [
        ("Engine ID: watchlist", 0),
        ("Flags transactions from known-suspicious supplier × country-of-origin pairs", 0),
        ("", 0),
        ("Algorithm:", 0),
        ("Look up (seller_id, seller_country) in the WATCHLIST set", 1),
        ("If present → flag the transaction", 1),
        ("If not → clear", 1),
        ("", 0),
        ("Simple binary lookup — no statistical analysis", 0),
        ("Editable in lib/watchlist.py, takes effect immediately", 0),
    ], font_size=11)

    _add_table(s3, 0.5, 3.5, 5.5, [
        ["Seller ID", "Country", "Supplier Name"],
        ["SUP001", "DE", "TechZone GmbH"],
        ["SUP002", "FR", "FashionHub Paris"],
        ["SUP005", "NL", "SportsPro Amsterdam"],
    ], col_widths=[1.5, 1.0, 3.0])

    # Note box
    txBox = s3.shapes.add_textbox(Inches(6.3), Inches(1.1), Inches(3.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Design rationale"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = EU_BLUE
    for text in [
        "",
        "Covers scenarios where intelligence identifies a supplier as high-risk regardless of their current VAT behaviour.",
        "",
        "Complements the statistical VAT ratio check with a rule-based layer.",
        "",
        "Adding new entries requires no code change — just edit the WATCHLIST set.",
    ]:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = text
        run2.font.size = Pt(10)
        run2.font.color.rgb = DARK_TEXT
        p2.space_after = Pt(2)

    # ── Slide 4: Risk Score Consolidation ──────────────────────────────
    s4 = _add_content_slide(prs, "Risk Score, Categories & Confidence")

    _add_bullet_box(s4, 0.5, 1.1, 4.8, 1.8, [
        ("Both engines publish to a single RT Risk Outcome broker", 0),
        ("The Automated Assessment Factory consolidates per transaction:", 0),
        ("", 0),
        ("Risk Score = flagged_count / total_outcomes_received", 1),
        ("If no outcomes received (timeout): score = 50% (uncertain)", 1),
        ("", 0),
        ("Confidence = outcomes_received / TOTAL_RISK_ENGINES", 1),
        ("With 2 engines: 0% (none), 50% (one), 100% (both)", 1),
    ], font_size=11)

    # Routing thresholds table
    _add_table(s4, 0.5, 3.2, 4.8, [
        ["Score Range", "Route", "Action"],
        ["< 33.33%", "Green → Release", "Auto-released, stored in DB"],
        ["33.33% – 66.66%", "Amber → Investigate", "Sent to C&T Risk Management"],
        ["> 66.66%", "Red → Retain", "Sent to C&T Risk Management"],
    ], col_widths=[1.3, 1.5, 2.0])

    # Effective mapping table (right side)
    _add_table(s4, 5.6, 1.1, 4.0, [
        ["Flagged", "Score", "Confidence", "Route"],
        ["0 of 2", "0%", "100%", "Release"],
        ["1 of 2", "50%", "100%", "Investigate"],
        ["2 of 2", "100%", "100%", "Retain"],
        ["0 of 1 (timeout)", "0%", "50%", "Release"],
        ["1 of 1 (timeout)", "100%", "50%", "Retain"],
        ["0 of 0 (both timeout)", "50%", "0%", "Investigate"],
    ], col_widths=[1.0, 0.8, 1.0, 1.2])

    # Timer note
    txBox = s4.shapes.add_textbox(Inches(5.6), Inches(3.5), Inches(4.0), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Assessment Timer (3 seconds)"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = EU_BLUE
    for text in [
        "",
        "Starts when Order Validation arrives.",
        "Publishes immediately if all engines respond early.",
        "Otherwise publishes on timer expiry with partial data.",
        "Late-arriving engine outcomes are discarded.",
    ]:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = text
        run2.font.size = Pt(10)
        run2.font.color.rgb = DARK_TEXT
        p2.space_after = Pt(2)

    # ── Save ──────────────────────────────────────────────────────────
    out = "docs/risk_monitoring_rules.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
