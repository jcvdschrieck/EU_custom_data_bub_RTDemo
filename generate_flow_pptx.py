#!/usr/bin/env python3
"""
Generate EU_Custom_DataHub_MessageFlow.pptx  (2 slides)

Slide 1 — Technical swimlane with colour legend:
  EU Custom Data Hub (top) · Ireland Investigation (bottom)

Slide 2 — Business & Functional process flow:
  Message journey from arrival at the EU Custom Data Hub to the Ireland
  Revenue investigation queue, presented without implementation detail.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn as _qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────────
EU_BLUE      = RGBColor(0x00, 0x33, 0x99)
EU_BLUE_LT   = RGBColor(0xD4, 0xDF, 0xF2)
EU_YELLOW    = RGBColor(0xFF, 0xED, 0x00)
IE_GREEN     = RGBColor(0x16, 0x9B, 0x62)
IE_GREEN_LT  = RGBColor(0xD0, 0xED, 0xDD)
ALARM_RED    = RGBColor(0xCC, 0x22, 0x22)
ALARM_RED_LT = RGBColor(0xFA, 0xE0, 0xE0)
AGENT_PURPLE = RGBColor(0x6A, 0x3D, 0x9A)
AGENT_PUR_LT = RGBColor(0xE8, 0xDD, 0xF5)
DB_TEAL      = RGBColor(0x00, 0x7A, 0x8A)
DB_TEAL_LT   = RGBColor(0xCC, 0xEA, 0xEE)
SIM_BLUE     = RGBColor(0x17, 0x6A, 0xA0)
SIM_BLUE_LT  = RGBColor(0xCE, 0xE5, 0xF2)
DECISION_AMB = RGBColor(0xE6, 0x8A, 0x00)   # amber for decision diamonds (slide 2)
DECISION_LT  = RGBColor(0xFD, 0xF0, 0xD8)
TERM_GREY    = RGBColor(0x6C, 0x75, 0x7D)   # terminal "no-action" boxes
TERM_GREY_LT = RGBColor(0xE9, 0xEC, 0xEF)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xF0, 0xF2, 0xF5)
MID_GREY     = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GREY    = RGBColor(0x22, 0x22, 0x22)


# ── Primitives ────────────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h,
          fill=WHITE, border=MID_GREY, border_pt=1.0,
          text="", font_size=10, bold=False,
          text_color=DARK_GREY, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE, wrap=True,
          shape_type=1):          # 1=rectangle, 4=diamond
    shp = slide.shapes.add_shape(shape_type,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_pt)
    if text:
        tf = shp.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = v_anchor
        p = tf.paragraphs[0]
        p.alignment = align
        for chunk in text.split("\n"):
            if p.runs:
                p = tf.add_paragraph()
                p.alignment = align
            run = p.add_run()
            run.text = chunk
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color
    return shp


def _txt(slide, text, x, y, w, h,
         font_size=10, bold=False, color=DARK_GREY,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for chunk in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = chunk
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _line(slide, x1, y1, x2, y2,
          color=DARK_GREY, width_pt=1.5,
          arrow_end=True, dashed=False):
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cx.line.color.rgb = color
    cx.line.width = Pt(width_pt)
    ln = cx.line._ln
    if dashed:
        prstDash = etree.SubElement(ln, _qn("a:prstDash"))
        prstDash.set("val", "dash")
    if arrow_end:
        tail = ln.find(_qn("a:tailEnd"))
        if tail is None:
            tail = etree.SubElement(ln, _qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return cx


def _pill(slide, x, y, w, h, label, fill, text_color=WHITE, font_size=7.5):
    _rect(slide, x, y, w, h, fill=fill, border=fill, border_pt=0,
          text=label, font_size=font_size, bold=True,
          text_color=text_color, align=PP_ALIGN.CENTER,
          v_anchor=MSO_ANCHOR.MIDDLE)


def _db_box(slide, x, y, w, h, label, fill=DB_TEAL, fill_lt=DB_TEAL_LT):
    _rect(slide, x, y, w, h, fill=fill, border=fill, border_pt=0,
          text=label, font_size=7.5, bold=True, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, x, y, w, h * 0.20, fill=fill_lt, border=fill, border_pt=0)


def _step_card(slide, x, y, w, h, step_num, title, body,
               hdr_fill, body_fill=WHITE, hdr_text=WHITE):
    HDR, PAD = 0.44, 0.10
    _rect(slide, x, y, w, h, fill=body_fill, border=hdr_fill, border_pt=1.5)
    _rect(slide, x, y, w, HDR, fill=hdr_fill, border=hdr_fill, border_pt=0)
    BADGE = 0.28
    bx, by = x + 0.09, y + (HDR - BADGE) / 2
    _rect(slide, bx, by, BADGE, BADGE, fill=WHITE, border=hdr_fill, border_pt=0,
          text=str(step_num), font_size=9, bold=True, text_color=hdr_fill,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, title, x + BADGE + 0.22, y + 0.06, w - BADGE - 0.32, HDR - 0.08,
         font_size=8.5, bold=True, color=hdr_text)
    _txt(slide, body, x + PAD, y + HDR + 0.09, w - PAD * 2, h - HDR - 0.12,
         font_size=7.5, color=DARK_GREY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Technical swimlane  (with colour legend)
# ══════════════════════════════════════════════════════════════════════════════

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_GREY

    # Title
    TITLE_H = 0.52
    _rect(slide, 0, 0, 13.33, TITLE_H, fill=EU_BLUE, border=EU_BLUE,
          text="EU Custom Data Hub — Technical Message Flow (Swimlane)",
          font_size=14, bold=True, text_color=WHITE,
          align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 0, 0, 0.14, TITLE_H, fill=EU_YELLOW, border=EU_YELLOW)

    # Lane geometry
    LANE_TOP  = TITLE_H + 0.06
    LABEL_W   = 1.18
    LANE_H    = (7.5 - LANE_TOP - 0.10) / 2
    EU_LANE_Y = LANE_TOP
    IE_LANE_Y = LANE_TOP + LANE_H

    _rect(slide, 0, EU_LANE_Y, 13.33, LANE_H, fill=EU_BLUE_LT, border=EU_BLUE_LT)
    _rect(slide, 0, IE_LANE_Y, 13.33, LANE_H, fill=IE_GREEN_LT, border=IE_GREEN_LT)
    _rect(slide, 0, IE_LANE_Y, 13.33, 0.03, fill=MID_GREY, border=MID_GREY)

    _rect(slide, 0, EU_LANE_Y, LABEL_W, LANE_H, fill=EU_BLUE, border=EU_BLUE,
          text="EU Custom\nData Hub\n\nEuropean\nCommission\nport 8505",
          font_size=8.5, bold=False, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 0, IE_LANE_Y, LABEL_W, LANE_H, fill=IE_GREEN, border=IE_GREEN,
          text="Ireland\nInvestigation\n\nIrish Revenue\nAgent + Queue",
          font_size=8.5, bold=False, text_color=WHITE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # Card geometry — 8 steps
    CONTENT_X = LABEL_W + 0.10
    CONTENT_W = 13.33 - CONTENT_X - 0.08
    N = 8; GAP = 0.09
    CARD_W  = (CONTENT_W - (N - 1) * GAP) / N
    PAD_V   = 0.16
    EU_CARD_Y = EU_LANE_Y + PAD_V
    IE_CARD_Y = IE_LANE_Y + PAD_V
    CARD_H    = LANE_H - 2 * PAD_V

    def cx(i): return CONTENT_X + i * (CARD_W + GAP)
    EU_MID = EU_LANE_Y + LANE_H / 2
    IE_MID = IE_LANE_Y + LANE_H / 2

    # ── Step cards ────────────────────────────────────────────────────────────
    _step_card(slide, cx(0), EU_CARD_Y, CARD_W, CARD_H, 1,
               "Simulation DB",
               "1,514 March 2026\ntransactions pre-seeded.\n\n"
               "Scenario week 2\n(8–14 Mar):\n"
               "TechZone GmbH → IE\n8 txn/day at 0%\n(correct: 23%)\n\n"
               "Replay speed:\nconfigurab. 1×–2880×",
               SIM_BLUE, SIM_BLUE_LT)
    _db_box(slide, cx(0)+0.09, EU_CARD_Y+CARD_H-0.44-0.08, CARD_W-0.18, 0.44,
            "simulation.db", SIM_BLUE, SIM_BLUE_LT)

    _step_card(slide, cx(1), EU_CARD_Y, CARD_W, CARD_H, 2,
               "Simulation Engine",
               "asyncio background\ntask; ticks every 50ms.\n\n"
               "Advances simulated\nclock at configured\nspeed.\n\n"
               "Fetches pending\ntransactions ≤ sim_time\nfrom simulation.db;\n"
               "marks fired=1.\n\n"
               "Calls fire_callback()",
               SIM_BLUE, SIM_BLUE_LT)

    _step_card(slide, cx(2), EU_CARD_Y, CARD_W, CARD_H, 3,
               "Insert to EU Custom DB",
               "insert_transaction()\nwrites each record to\neuropean_custom.db.\n\n"
               "Transaction added to\nin-memory live queue\n(30-tx ring buffer\n"
               "for front-end poll).\n\n"
               "All records stored\nregardless of VAT\ncorrectness.",
               EU_BLUE, WHITE)
    _db_box(slide, cx(2)+0.09, EU_CARD_Y+CARD_H-0.44-0.08, CARD_W-0.18, 0.44,
            "european_custom.db")

    _step_card(slide, cx(3), EU_CARD_Y, CARD_W, CARD_H, 4,
               "VAT Alarm Checker",
               "After DB write:\n\n"
               "7-day VAT/value ratio\nvs 8-week baseline.\n\n"
               "Deviation > 25%:\n→ alarm raised\n   (7-day validity)\n\n"
               "Active alarm + IE:\n→ suspicious = 1\n   suspicion: MEDIUM\n\n"
               "Non-IE: alarm only,\nno suspicious flag.",
               ALARM_RED, ALARM_RED_LT)
    _pill(slide, cx(3)+0.09, EU_CARD_Y+CARD_H-0.34, CARD_W-0.18, 0.26,
          "⚠ alarms table", ALARM_RED)

    _step_card(slide, cx(4), EU_CARD_Y, CARD_W, CARD_H, 5,
               "Agent Queue",
               "If suspicious=1\n(IE-bound only):\ntx enqueued to\n"
               "asyncio.Queue\n(non-blocking).\n\n"
               "Background worker\npicks up items;\nruns agent in\n"
               "ThreadPoolExecutor\n(keeps sim loop free).\n\n"
               "Queue depth shown\nin simulation/status.",
               AGENT_PURPLE, AGENT_PUR_LT)

    _step_card(slide, cx(5), IE_CARD_Y, CARD_W, CARD_H, 6,
               "VAT Fraud Detection Agent",
               "Subprocess bridge:\nagent_bridge.py calls\n_analyse_tx.py in\n"
               "vat_fraud_detection/\n\n"
               "Builds Invoice +\nLineItem from tx dict.\n\n"
               "RAG retrieves Irish\nVAT legislation from\nChromaDB.\n\n"
               "LM Studio LLM returns\nverdict per line item:\ncorrect / incorrect\n/ uncertain",
               AGENT_PURPLE, AGENT_PUR_LT)
    _pill(slide, cx(5)+0.09, IE_CARD_Y+CARD_H-0.34, CARD_W-0.18, 0.26,
          "LM Studio  port 1234", AGENT_PURPLE)

    _step_card(slide, cx(6), IE_CARD_Y, CARD_W, CARD_H, 7,
               "Verdict Routing",
               "INCORRECT:\n"
               "→ suspicion_level\n   upgraded to HIGH\n"
               "→ insert_ireland_\n   queue()\n"
               "→ insert_agent_\n   log(sent=1)\n\n"
               "CORRECT /\nUNCERTAIN:\n"
               "→ suspicious flag\n   cleared\n"
               "→ insert_agent_\n   log(sent=0)",
               IE_GREEN, IE_GREEN_LT)
    _pill(slide, cx(6)+0.09, IE_CARD_Y+CARD_H-0.62, CARD_W-0.18, 0.24,
          "INCORRECT → HIGH", ALARM_RED)
    _pill(slide, cx(6)+0.09, IE_CARD_Y+CARD_H-0.34, CARD_W-0.18, 0.24,
          "CORRECT / UNCERTAIN → cleared", IE_GREEN)

    _step_card(slide, cx(7), IE_CARD_Y, CARD_W, CARD_H, 8,
               "Ireland Investigation Queue",
               "Confirmed-incorrect\ncases for local\ninvestigation.\n\n"
               "Each entry includes:\n· Transaction detail\n"
               "· Alarm key + deviation\n"
               "· Applied vs correct\n  VAT rate\n"
               "· Agent reasoning\n· Suspicion: HIGH\n\n"
               "Front-end: /ireland",
               IE_GREEN, IE_GREEN_LT)
    _db_box(slide, cx(7)+0.09, IE_CARD_Y+CARD_H-0.44-0.08, CARD_W-0.18, 0.44,
            "ireland_queue table", IE_GREEN, IE_GREEN_LT)

    # ── Within-lane arrows ────────────────────────────────────────────────────
    for i in range(4): _line(slide, cx(i)+CARD_W, EU_MID, cx(i+1), EU_MID,
                              color=EU_BLUE, width_pt=2.0)
    for i in [5, 6]:  _line(slide, cx(i)+CARD_W, IE_MID, cx(i+1), IE_MID,
                              color=IE_GREEN, width_pt=2.0)

    # ── Cross-lane arrow: Agent Queue (step 5) → Agent (step 6) ──────────────
    LANE_SEP    = IE_LANE_Y
    EU_CARD_BTM = EU_CARD_Y + CARD_H
    IE_CARD_TOP = IE_CARD_Y
    X5 = cx(4) + CARD_W * 0.55
    X6 = cx(5) + CARD_W * 0.45
    _line(slide, X5, EU_CARD_BTM, X5, LANE_SEP,   color=AGENT_PURPLE, width_pt=2.0, arrow_end=False)
    _line(slide, X5, LANE_SEP,    X6, LANE_SEP,   color=AGENT_PURPLE, width_pt=2.0, arrow_end=False)
    _line(slide, X6, LANE_SEP,    X6, IE_CARD_TOP, color=AGENT_PURPLE, width_pt=2.0, arrow_end=True)
    _txt(slide, "suspicious IE tx\n(suspicion: MEDIUM)",
         min(X5,X6)-0.05, LANE_SEP-0.38, 1.30, 0.38,
         font_size=7, color=AGENT_PURPLE, italic=True, align=PP_ALIGN.CENTER)

    # ── Colour legend ─────────────────────────────────────────────────────────
    LEG_X = 0.18
    LEG_Y = 7.5 - 0.54
    LEG_H = 0.46
    LEG_W_TOT = 13.0

    _rect(slide, LEG_X-0.08, LEG_Y-0.04, LEG_W_TOT+0.16, LEG_H+0.08,
          fill=WHITE, border=MID_GREY, border_pt=0.5)

    _txt(slide, "COLOUR CODE", LEG_X, LEG_Y, 1.0, 0.18,
         font_size=7, bold=True, color=DARK_GREY)

    items = [
        (SIM_BLUE,     "Simulation layer (replay engine)"),
        (EU_BLUE,      "EU Custom Data Hub processes"),
        (ALARM_RED,    "VAT alarm detection"),
        (AGENT_PURPLE, "AI agent processing"),
        (IE_GREEN,     "Ireland investigation"),
        (DB_TEAL,      "Persistent database storage"),
        (AGENT_PURPLE, "Cross-lane handoff (suspicious IE tx)"),
    ]
    pill_w, pill_h = 0.18, 0.16
    lbl_w = 1.52
    item_w = pill_w + lbl_w + 0.08
    x_cur = LEG_X
    for fill, label in items:
        _pill(slide, x_cur, LEG_Y+0.22, pill_w, pill_h, "", fill)
        _txt(slide, label, x_cur+pill_w+0.04, LEG_Y+0.20, lbl_w, 0.20,
             font_size=6.5, color=DARK_GREY)
        x_cur += item_w


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Business & Functional process flow
# ══════════════════════════════════════════════════════════════════════════════

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_GREY

    # ── Title ─────────────────────────────────────────────────────────────────
    TITLE_H = 0.52
    _rect(slide, 0, 0, 13.33, TITLE_H, fill=EU_BLUE, border=EU_BLUE,
          text="Transaction Journey — Business & Functional Perspective",
          font_size=14, bold=True, text_color=WHITE,
          align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 0, 0, 0.14, TITLE_H, fill=EU_YELLOW, border=EU_YELLOW)

    # ── Context banner ────────────────────────────────────────────────────────
    _rect(slide, 0.15, 0.58, 13.03, 0.34, fill=EU_BLUE_LT, border=EU_BLUE,
          border_pt=0.8,
          text="Cross-border B2C e-commerce transaction (supplier → EU buyer) "
               "arrives at the EU Custom Data Hub carrying: seller identity, buyer country, "
               "product category, value, and applied VAT rate.",
          font_size=8.5, bold=False, text_color=EU_BLUE,
          align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── Layout constants ──────────────────────────────────────────────────────
    # Main flow row: process boxes + decision diamonds
    Y_MID   = 2.64      # vertical centre of main flow elements
    BOX_H   = 1.22      # height of process boxes
    BOX_W   = 1.32      # width of process boxes
    DIA_H   = 1.10      # diamond height
    DIA_W   = 1.14      # diamond width
    ARR_W   = 0.22      # horizontal gap used by arrows
    Y_BOX_T = Y_MID - BOX_H / 2
    Y_DIA_T = Y_MID - DIA_H / 2
    Y_BOX_B = Y_MID + BOX_H / 2
    Y_DIA_B = Y_MID + DIA_H / 2

    # Branch boxes below decision diamonds
    Y_BRANCH_MID = 5.30
    BRANCH_H     = 0.90
    BRANCH_W     = 1.46
    Y_BRANCH_T   = Y_BRANCH_MID - BRANCH_H / 2

    # X positions of every element (left edge), computed sequentially
    MARGIN = 0.12
    def _seq(*widths):
        """Return list of left-edge x positions given alternating widths/gaps."""
        xs, cur = [], MARGIN
        for w in widths:
            xs.append(cur); cur += w
        return xs

    # Elements: B=box, D=diamond, A=arrow-gap
    # 6 boxes, 3 diamonds, arrows between each pair
    A = ARR_W
    xs = _seq(
        BOX_W, A,    # 0: "Transaction Arrives"           | arrow
        BOX_W, A,    # 2: "Transaction Recorded"          | arrow
        BOX_W, A,    # 4: "VAT Compliance Monitoring"     | arrow
        DIA_W, A,    # 6: ◇ "Deviation alert?"            | arrow YES
        BOX_W, A,    # 8: "Alert Raised"                  | arrow
        DIA_W, A,    # 10: ◇ "Ireland-bound?"             | arrow YES
        BOX_W, A,    # 12: "AI Compliance Agent"          | arrow
        DIA_W, A,    # 14: ◇ "VAT confirmed incorrect?"   | arrow YES
        BOX_W,       # 16: "Ireland Revenue Queue"
    )
    # xs indices: 0,2,4,6,8,10,12,14,16 = element left edges; 1,3,5,7,9,11,13,15 = arrow gaps

    X = {
        'arrive':   xs[0],   # process box 1
        'record':   xs[2],   # process box 2
        'monitor':  xs[4],   # process box 3
        'alarm_d':  xs[6],   # diamond 1
        'alert':    xs[8],   # process box 4
        'ie_d':     xs[10],  # diamond 2
        'agent':    xs[12],  # process box 5
        'verd_d':   xs[14],  # diamond 3
        'ireland':  xs[16],  # process box 6
    }

    # Helper: centre x of an element given its left edge and width
    def mid_x(key, w): return X[key] + w / 2

    # ── Process boxes (main row) ──────────────────────────────────────────────

    # Box 1: Transaction arrives
    _rect(slide, X['arrive'], Y_BOX_T, BOX_W, BOX_H,
          fill=EU_BLUE_LT, border=EU_BLUE, border_pt=2.0,
          text="Transaction\nArrives\nat EU Custom\nData Hub",
          font_size=9.5, bold=True, text_color=EU_BLUE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    # Actor tag
    _pill(slide, X['arrive'], Y_BOX_T-0.26, BOX_W, 0.22,
          "Supplier (cross-border B2C)", EU_BLUE, font_size=6.5)

    # Box 2: Transaction recorded
    _rect(slide, X['record'], Y_BOX_T, BOX_W, BOX_H,
          fill=EU_BLUE_LT, border=EU_BLUE, border_pt=2.0,
          text="Transaction\nRecorded\nPermanently",
          font_size=9.5, bold=True, text_color=EU_BLUE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, "Regardless of\ncompliance status",
         X['record']+0.08, Y_BOX_T+BOX_H*0.62, BOX_W-0.16, 0.30,
         font_size=7, color=EU_BLUE, align=PP_ALIGN.CENTER, italic=True)
    _pill(slide, X['record'], Y_BOX_T-0.26, BOX_W, 0.22,
          "EU Custom Data Hub", EU_BLUE, font_size=6.5)

    # Box 3: VAT compliance monitoring
    _rect(slide, X['monitor'], Y_BOX_T, BOX_W, BOX_H,
          fill=EU_BLUE_LT, border=EU_BLUE, border_pt=2.0,
          text="VAT Compliance\nMonitoring\n(automated)",
          font_size=9.5, bold=True, text_color=EU_BLUE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, "7-day ratio vs\n8-week baseline",
         X['monitor']+0.08, Y_BOX_T+BOX_H*0.62, BOX_W-0.16, 0.30,
         font_size=7, color=EU_BLUE, align=PP_ALIGN.CENTER, italic=True)
    _pill(slide, X['monitor'], Y_BOX_T-0.26, BOX_W, 0.22,
          "EU Custom Data Hub", EU_BLUE, font_size=6.5)

    # Diamond 1: Alarm?
    _rect(slide, X['alarm_d'], Y_DIA_T, DIA_W, DIA_H,
          fill=DECISION_LT, border=DECISION_AMB, border_pt=2.0,
          text="VAT Rate\nDeviation\nAlert?",
          font_size=8.5, bold=True, text_color=DECISION_AMB,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
          shape_type=4)   # diamond

    # Box 4: Alert raised
    _rect(slide, X['alert'], Y_BOX_T, BOX_W, BOX_H,
          fill=ALARM_RED_LT, border=ALARM_RED, border_pt=2.0,
          text="Deviation Alert\nRaised\n(7-day validity)",
          font_size=9.5, bold=True, text_color=ALARM_RED,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, "Supplier flagged for\nenhanced monitoring",
         X['alert']+0.08, Y_BOX_T+BOX_H*0.62, BOX_W-0.16, 0.30,
         font_size=7, color=ALARM_RED, align=PP_ALIGN.CENTER, italic=True)
    _pill(slide, X['alert'], Y_BOX_T+BOX_H+0.04, BOX_W, 0.22,
          "Suspicion: MEDIUM", ALARM_RED, font_size=6.5)

    # Diamond 2: Ireland-bound?
    _rect(slide, X['ie_d'], Y_DIA_T, DIA_W, DIA_H,
          fill=DECISION_LT, border=DECISION_AMB, border_pt=2.0,
          text="Destination:\nIreland?",
          font_size=8.5, bold=True, text_color=DECISION_AMB,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
          shape_type=4)

    # Box 5: AI compliance agent
    _rect(slide, X['agent'], Y_BOX_T, BOX_W, BOX_H,
          fill=AGENT_PUR_LT, border=AGENT_PURPLE, border_pt=2.0,
          text="AI Compliance\nAgent\nAnalysis",
          font_size=9.5, bold=True, text_color=AGENT_PURPLE,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, "VAT legislation RAG\n+ LLM verdict\nper line item",
         X['agent']+0.08, Y_BOX_T+BOX_H*0.56, BOX_W-0.16, 0.38,
         font_size=7, color=AGENT_PURPLE, align=PP_ALIGN.CENTER, italic=True)
    _pill(slide, X['agent'], Y_BOX_T-0.26, BOX_W, 0.22,
          "EU Customs → Irish Revenue Agent", AGENT_PURPLE, font_size=6.5)

    # Diamond 3: Verdict incorrect?
    _rect(slide, X['verd_d'], Y_DIA_T, DIA_W, DIA_H,
          fill=DECISION_LT, border=DECISION_AMB, border_pt=2.0,
          text="VAT Rate\nConfirmed\nIncorrect?",
          font_size=8.5, bold=True, text_color=DECISION_AMB,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE,
          shape_type=4)

    # Box 6: Ireland Revenue investigation queue
    _rect(slide, X['ireland'], Y_BOX_T, BOX_W, BOX_H,
          fill=IE_GREEN_LT, border=IE_GREEN, border_pt=2.5,
          text="Ireland Revenue\nInvestigation\nQueue",
          font_size=9.5, bold=True, text_color=IE_GREEN,
          align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, "Alarm evidence +\nagent reasoning\nprovided",
         X['ireland']+0.08, Y_BOX_T+BOX_H*0.58, BOX_W-0.16, 0.36,
         font_size=7, color=IE_GREEN, align=PP_ALIGN.CENTER, italic=True)
    _pill(slide, X['ireland'], Y_BOX_T-0.26, BOX_W, 0.22,
          "Irish Revenue Authority", IE_GREEN, font_size=6.5)
    _pill(slide, X['ireland'], Y_BOX_T+BOX_H+0.04, BOX_W, 0.22,
          "Suspicion: HIGH", ALARM_RED, font_size=6.5)

    # ── "No action" terminal states (below each decision diamond) ─────────────

    # Below diamond 1 ("No deviation"): X centre = alarm_d centre
    D1_CX = X['alarm_d'] + DIA_W / 2
    D2_CX = X['ie_d']    + DIA_W / 2
    D3_CX = X['verd_d']  + DIA_W / 2

    def _terminal(cx, label, sub, fill=TERM_GREY, fill_lt=TERM_GREY_LT):
        bw = BRANCH_W
        bx = cx - bw / 2
        _rect(slide, bx, Y_BRANCH_T, bw, BRANCH_H,
              fill=fill_lt, border=fill, border_pt=1.5,
              text=label, font_size=8, bold=True, text_color=fill,
              align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        _txt(slide, sub, bx+0.08, Y_BRANCH_T+BRANCH_H*0.56, bw-0.16, 0.30,
             font_size=7, color=fill, align=PP_ALIGN.CENTER, italic=True)

    _terminal(D1_CX,
              "No Further Action",
              "Transaction within\nnormal VAT parameters")

    _terminal(D2_CX,
              "Under Monitoring",
              "Alert recorded;\nno member-state referral")

    _terminal(D3_CX,
              "Suspicion Cleared",
              "Correct or inconclusive;\nno further investigation")

    # ── Main-flow horizontal arrows ───────────────────────────────────────────
    pairs = [
        ('arrive',  BOX_W, 'record',  BOX_W),
        ('record',  BOX_W, 'monitor', BOX_W),
        ('monitor', BOX_W, 'alarm_d', DIA_W),
        ('alarm_d', DIA_W, 'alert',   BOX_W),
        ('alert',   BOX_W, 'ie_d',    DIA_W),
        ('ie_d',    DIA_W, 'agent',   BOX_W),
        ('agent',   BOX_W, 'verd_d',  DIA_W),
        ('verd_d',  DIA_W, 'ireland', BOX_W),
    ]
    for src, sw, dst, _ in pairs:
        x_from = X[src] + sw
        x_to   = X[dst]
        _line(slide, x_from, Y_MID, x_to, Y_MID,
              color=EU_BLUE, width_pt=2.0)

    # YES label on diamond outgoing arrows
    yes_diamonds = [
        (X['alarm_d'] + DIA_W + 0.01, Y_MID - 0.22, "YES"),
        (X['ie_d']    + DIA_W + 0.01, Y_MID - 0.22, "YES"),
        (X['verd_d']  + DIA_W + 0.01, Y_MID - 0.22, "YES"),
    ]
    for lx, ly, lbl in yes_diamonds:
        _txt(slide, lbl, lx, ly, 0.20, 0.20,
             font_size=7.5, bold=True, color=IE_GREEN, align=PP_ALIGN.CENTER)

    # ── Downward "NO" arrows from each diamond ────────────────────────────────
    for dcx, dbt in [(D1_CX, Y_DIA_B), (D2_CX, Y_DIA_B), (D3_CX, Y_DIA_B)]:
        _line(slide, dcx, dbt, dcx, Y_BRANCH_T,
              color=ALARM_RED, width_pt=1.8, dashed=True)
        _txt(slide, "NO", dcx+0.04, dbt+0.06, 0.24, 0.18,
             font_size=7.5, bold=True, color=ALARM_RED, align=PP_ALIGN.LEFT)

    # ── Colour legend ─────────────────────────────────────────────────────────
    LEG_Y  = 7.5 - 0.50
    LEG_X  = 0.18
    _rect(slide, LEG_X-0.08, LEG_Y-0.04, 13.16, 0.48,
          fill=WHITE, border=MID_GREY, border_pt=0.5)
    _txt(slide, "COLOUR CODE", LEG_X, LEG_Y, 1.0, 0.18,
         font_size=7, bold=True, color=DARK_GREY)

    legend_items = [
        (EU_BLUE,      "EU Custom Data Hub process"),
        (ALARM_RED,    "VAT deviation alert"),
        (DECISION_AMB, "Decision / routing point"),
        (AGENT_PURPLE, "AI compliance agent"),
        (IE_GREEN,     "Ireland Revenue investigation"),
        (TERM_GREY,    "Terminal — no further action"),
    ]
    pill_w, pill_h = 0.18, 0.16
    lbl_w  = 1.60
    x_cur  = LEG_X
    for fill, label in legend_items:
        _pill(slide, x_cur, LEG_Y+0.24, pill_w, pill_h, "", fill)
        _txt(slide, label, x_cur+pill_w+0.04, LEG_Y+0.22, lbl_w, 0.20,
             font_size=6.5, color=DARK_GREY)
        x_cur += pill_w + lbl_w + 0.14


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    build_slide1(prs)
    build_slide2(prs)

    out = "EU_Custom_DataHub_MessageFlow.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
